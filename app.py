from dotenv import load_dotenv
load_dotenv()
import os
import re
import streamlit as st
import textwrap
import google.generativeai as genai
from pptx.util import Pt
from pptx.util import Inches
from pptx import Presentation

from tempfile import NamedTemporaryFile
from io import StringIO
from manim_gen import *

genai.configure(api_key=os.getenv('GOOGLE_API_KEY'))
model = genai.GenerativeModel('gemini-1.5-pro-latest')

#Asks Gemini to read a lesson plan PDF and create a pseudo-JSON slides outline

prompt = """Put all of the contents of the lesson plan into a detailed template for a slide show in JSON format.
Take following json as an input, generate elements content based on prompt and prepare and return json structure by combining all elements such as: {'heading': '<value>'}.

We are trying to create slide shows for a lesson plan Fill in the following elements with content from the lesson plan to create as many slides as needed to display all the content.
Create one overarching intro slide and name it. For the following slides add a title for each slide (appx 40 chars), and 1-2 sentence teaching the content for each slide , add a description of the math examples is being taught or add an example problem if not provided. Replace <value> in plain text.

{
  "topic title": "<value>",
  "subtext": "<value>",
  "example": "<value>"
}

Let the json file be titled "slides". Make sure the title slide has a key "example", but the value must be empty."""

content = []
#Since Gemini doesn't actually return a JSON, manually parses the response (oh well)

def load_content(response):
  text = response.text
    #Removes extra spaces, tabs, and new lines

  text = re.sub('([\t\n]+|\s{2,})', '', text)
    #Removes bracketing, quotes, and JSON padding

  text = re.sub('(\},|\",|"|[{}[[\]]|```json|```|slides: )', '', text)
    #Removes dictionary keys and only keep values

  text = re.sub('(heading: |subtext: |example: )', '\n', text)
  #Indexes content array by new lines

  content = re.split('\n', text)
  return content

pptx = Presentation('content/PresTemplate.pptx')
transparent_image = 'content/Transparent_Square_Tiles_Texture.png'

video_file_paths = []

st.set_page_config(page_title = "Test Presentation Generator")
st.header("Test Presentation Generator")
upload = st.file_uploader("Upload Lesson Plan PDF: ", key = "upload")


#Creates PowerPoint

def generate_pptx(pptx, refined_content):
  for i in range(1, len(refined_content), 3):
    layout = pptx.slide_layouts[7]
    slide = pptx.slides.add_slide(layout).shapes
    title = slide.title
    title.text = refined_content[i]
    title.text_frame.paragraphs[0].font.size = Pt(36)

    body = slide.placeholders[2]
    subtext = body.text_frame.add_paragraph()
    subtext.text = refined_content[i + 1]
    subtext.font.size = Pt(18)
    # subtext.space_after = Pt(4)
    # example = body.text_frame.add_paragraph()
    # example.text = refined_content[i + 2]
    # example.font.size = Pt(18)

    concept = refined_content[i + 1] + ' ' + refined_content[i + 2]

    # Skips title slide
    
    try:
        if (i > 1):
          #Creates Manim animations given the subtext and example and embeds the animation into the slide
          concept = refined_content[i + 1] + ' ' + refined_content[i + 2]
          video_file = make_animation(concept)
          video = slide.add_movie(video_file, Inches(6), Inches(2), Inches(5), Inches(2.8125), transparent_image, 'video/mp4')
          
    except exception as e:
        print(e)
        pass


  xml_slides = pptx.slides._sldIdLst
  slides = list(xml_slides)
  xml_slides.remove(slides[0])
  return pptx


#Creates a download button for the file

def download(file_path, topic):
  content = open(file_path, "rb").read()

  st.download_button(
      label = "Download",
      data = content,
      file_name = f"{topic}.pptx",
      key = "download_button"
  )


#Once uploaded, proceeds

if upload is not None:
  file_path = os.path.join("/tmp", upload.name)
  with open(file_path, "wb") as f:
      f.write(upload.getbuffer())

  uploaded_file = genai.upload_file(file_path, mime_type = 'application/pdf')
  response = model.generate_content([uploaded_file, prompt], request_options = {"timeout": 600})

  topic = st.text_input("File Name: ", key = "topic")
  submitted = st.button("Create Presentation")

  content = load_content(response)


  #Once submitted, calls Gemini and saves the resulting powerpoint

  if submitted:
    content = load_content(response)
    pptx = generate_pptx(pptx, content)

    pptx.save(f'{topic}.pptx')

    st.text("Ready to Download")
    download(f'{topic}.pptx', topic)